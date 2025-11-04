from dotenv import load_dotenv
import os
import re
import json
import html
from openai import OpenAI
import requests
from fastapi import Request, FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field
import cohere
from fastapi import Body
from fastapi import Request
from fastapi import FastAPI, UploadFile, File
from fastapi.responses import JSONResponse
from google.cloud import vision
from fastapi.responses import PlainTextResponse
from fastapi.middleware.cors import CORSMiddleware
from PyPDF2 import PdfReader
import tempfile
from google.oauth2 import service_account
from google.cloud import translate_v2 as google_translate
from google.cloud import vision
import html
import time
import io
import os
import re
import zipfile
import tempfile
import pathlib
import subprocess
import xml.etree.ElementTree as ET
import speech_recognition as sr  # 음성인식
from pydub import AudioSegment
from io import BytesIO
import imageio_ffmpeg
from hanspell import spell_checker
from datetime import datetime
from typing import Optional, Tuple, List
import fitz
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
import olefile
import chardet

try:
    import kss
    HAS_KSS = True
except Exception:
    HAS_KSS = False

_KO_TERMINALS = ("다.", "요.", "입니다.", "인가요?", "일까요?",
                 "습니까?", "했다.", "했다고", "했다며")

def split_ko_sentences(text: str):
    text = (text or "").strip()
    if not text:
        return []
    if HAS_KSS:
        # kss는 줄바꿈 포함해도 잘 쪼개줌
        return [s.strip() for s in kss.split_sentences(text) if s.strip()]
    # fallback: 아주 단순한 규칙
    #  - 줄바꿈을 공백으로 바꾸고
    #  - 문장말(다./요./입니다./?! 등) 기준 대략 분리
    t = re.sub(r"[ \t]*\n[ \t]*", " ", text)
    parts = re.split(r"(?<=[\.!?]|다\.|요\.|입니다\.)\s+", t)
    return [p.strip() for p in parts if p.strip()]

def count_ko_sentences(text: str) -> int:
    return len(split_ko_sentences(text))

def _expand_clamp(v, lo, hi):
    try:
        v = int(v)
    except Exception:
        v = lo
    return max(lo, min(v, hi))

def _expand_split_sentences_kor(text: str):
    """한글 문장 경계 대충 잘라내는 경량 분리기(마침표/요/니다/… 기준)."""
    t = re.sub(r'\s*\n+\s*', '\n', (text or '').strip())
    s = re.sub(r'([\.!?…])\s+', r'\1\n', t)
    s = re.sub(r'(다\.|요\.|니다\.)\s+', r'\1\n', s)
    parts = [p.strip() for p in s.split('\n') if p.strip()]
    return parts

def _expand_ends_with_terminal(s: str) -> bool:
    return bool(re.search(r'(다\.|요\.|니다\.)\s*$|[.!?…]\s*$|[”’"\')\]】〉》」』]\s*$', s or ''))

def _expand_crop_to_last_boundary(s: str) -> str:
    """마지막 문장 경계(마침표/따옴표 등)에서 컷."""
    if not s:
        return s
    m = list(re.finditer(r'(다\.|요\.|니다\.)|[.!?…]|[”’"\')\]】〉》」』]', s))
    return s if not m else s[:m[-1].end()]

def _expand_smart_trim_to_chars(text: str, max_chars: int) -> str:
    """문장 경계 유지하며 글자 상한에 맞춤. 실패 시 하드 컷."""
    if len(text) <= max_chars:
        return text
    parts = _expand_split_sentences_kor(text)
    out, total = [], 0
    for p in parts:
        if not _expand_ends_with_terminal(p):
            continue
        add = len(p) + (1 if out else 0)
        if total + add > max_chars:
            break
        out.append(p)
        total += add
    return ("\n".join(out).rstrip() if out else text[:max_chars].rstrip())

def _expand_est_tokens_for_chars(chars: int) -> int:
    """대략적인 토큰 예산(한글 기준 여유있게)."""
    return max(128, int(chars / 1.6) + 96)

def _expand_llm(client, messages, max_chars: int, temperature=0.35):
    """Chat Completions 래퍼: max_tokens를 글자수 기반으로 안전 설정."""
    max_tokens = _expand_est_tokens_for_chars(max_chars)
    resp = client.chat.completions.create(
        model="gpt-4",                 # ← 필요시 사용 중인 모델명으로 통일
        messages=messages,
        temperature=temperature,
        max_tokens=max_tokens,
    )
    choice = resp.choices[0]
    # 일부 SDK에서 finish_reason 접근 방법이 달라질 수 있으므로 getattr 사용
    finish = getattr(choice, "finish_reason",
                     getattr(choice, "finish_reason", None))
    text = (choice.message.content or "").strip()
    return text, finish

def _expand_count_sents(s: str) -> int:
    return len(_expand_split_sentences_kor(s or ""))


class InformalInput(BaseModel):
    content: str
    ending: Optional[str] = None  # 'hada' | 'haetda' | 'hae'

class ExpandInput(BaseModel):
    content: str
    length_boost: int = 20
    length_level: Optional[str] = None
    add_sentences: int = 0

    class Config:
        extra = "ignore"

def _clamp(v, lo, hi):
    return max(lo, min(v, hi))

def _split_sentences_kor(text: str):
    """한국어/영문 혼용 간단 분리: 마침표류/종결형 + 개행 기준."""
    text = re.sub(r'\s*\n+\s*', '\n', text.strip())
    s = re.sub(r'([\.!?…])\s+', r'\1\n', text)
    s = re.sub(r'(다\.|요\.|니다\.)\s+', r'\1\n', s)
    return [p.strip() for p in s.split('\n') if p.strip()]

def _smart_trim_to_chars(text: str, max_chars: int) -> str:
    """문장 단위로 합치면서 상한 이하로 유지(불완전 마지막 문장은 버림)."""
    if len(text) <= max_chars:
        return text
    parts = _split_sentences_kor(text)
    out, total = [], 0
    for p in parts:
        # 완결 문장만 추가
        if not _ends_with_terminal(p):
            continue
        if total + len(p) + (1 if out else 0) > max_chars:
            break
        out.append(p)
        total += len(p) + (1 if out else 0)
    # 여전히 비면 하드컷
    return ("\n".join(out).rstrip() if out else text[:max_chars].rstrip())

def _ends_with_terminal(s: str) -> bool:
    """문장 완결 판단(한국어 종결/문장부호/닫힘문자)."""
    return bool(re.search(r'(다\.|요\.|니다\.)\s*$|[.!?…]\s*$|[”’"\')\]】〉》」』]\s*$', s))

def _crop_to_last_boundary(s: str) -> str:
    """마지막 문장 경계까지만 남김."""
    m = list(re.finditer(r'(다\.|요\.|니다\.)|[.!?…]|[”’"\')\]】〉》」』]', s))
    return s if not m else s[:m[-1].end()]

# 문서 이미지 텍스트 추출 관련
def ocr_image_bytes(img_bytes: bytes) -> str:
    if not img_bytes:
        return ""
    # ✅ 가드 추가: 초기화 실패시 바로 빈 문자열, 로그 남김
    if vision_client is None:
        print("[OCR] vision_client is None (check GOOGLE_APPLICATION_CREDENTIALS_JSON or GOOGLE_APPLICATION_CREDENTIALS)")
        return ""
    try:
        image = vision.Image(content=img_bytes)
        resp = vision_client.text_detection(image=image)
        raw = ""
        if getattr(resp, "full_text_annotation", None):
            raw = resp.full_text_annotation.text or ""
        elif resp.text_annotations:
            raw = resp.text_annotations[0].description or ""
        return clean_ocr_text(raw)
    except Exception as e:
        print("⚠️ OCR 실패:", e)
        return ""

# ---------------------------
# 포맷별 추출기 (본문 + 내장 이미지)
# 각 함수는 (본문텍스트, 이미지OCR리스트) 튜플 반환
# ---------------------------
def extract_from_image(raw: bytes) -> Tuple[str, List[str]]:
    # 단일 이미지 파일 자체 OCR
    return "", [ocr_image_bytes(raw)]

def extract_from_txt(raw: bytes) -> Tuple[str, List[str]]:
    return detect_text_encoding_and_decode(raw).strip(), []

def extract_from_pdf(raw: bytes) -> Tuple[str, List[str]]:
    text_parts, ocr_parts = [], []
    doc = fitz.open(stream=raw, filetype="pdf")
    for page in doc:
        # 본문 텍스트
        text_parts.append(page.get_text())
        # 페이지 내 이미지 → OCR
        for img in page.get_images(full=True):
            xref = img[0]
            try:
                meta = doc.extract_image(xref)
                img_bytes = meta.get("image", b"")
                if img_bytes:
                    ocr_parts.append(ocr_image_bytes(img_bytes))
            except Exception as e:
                print("pdf image extract err:", e)
    return "\n".join(text_parts).strip(), [t for t in ocr_parts if t]

def extract_from_docx(raw: bytes) -> Tuple[str, List[str]]:
    import zipfile
    import io
    text_parts, ocr_parts = [], []
    # 1) 본문 텍스트
    try:
        d = Document(io.BytesIO(raw))
        for p in d.paragraphs:
            text_parts.append(p.text or "")
    except Exception as e:
        print("docx text parse err:", e)
    # 2) 내장 이미지 (ZIP 스캔)
    try:
        found = 0
        with zipfile.ZipFile(io.BytesIO(raw)) as zf:
            for name in zf.namelist():
                low = name.lower()
                if low.startswith("word/media/") and low.split(".")[-1] in (
                    "png", "jpg", "jpeg", "bmp", "gif", "tif", "tiff", "webp"
                ):
                    ocr_parts.append(ocr_image_bytes(zf.read(name)))
                    found += 1
                print(
                    f"[DOCX] media images found: {found}, ocred: {sum(1 for t in ocr_parts if t)}")
    except Exception as e:
        print("docx media zip scan err:", e)
    return "\n".join(text_parts).strip(), [t for t in ocr_parts if t]

def extract_from_pptx(raw: bytes) -> Tuple[str, List[str]]:
    import zipfile
    import io
    text_parts, ocr_parts = [], []
    # 1) 본문 텍스트
    try:
        prs = Presentation(io.BytesIO(raw))
        for slide in prs.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text_parts.append(shape.text or "")
    except Exception as e:
        print("pptx text parse err:", e)
    # 2) 내장 이미지 (ZIP 스캔)
    try:
        with zipfile.ZipFile(io.BytesIO(raw)) as zf:
            for name in zf.namelist():
                low = name.lower()
                if low.startswith("ppt/media/") and low.split(".")[-1] in (
                    "png", "jpg", "jpeg", "bmp", "gif", "tif", "tiff", "webp"
                ):
                    ocr_parts.append(ocr_image_bytes(zf.read(name)))
    except Exception as e:
        print("pptx media zip scan err:", e)
    return "\n".join(text_parts).strip(), [t for t in ocr_parts if t]

def extract_from_xlsx(raw: bytes) -> Tuple[str, List[str]]:
    import zipfile
    import io
    text_parts, ocr_parts = [], []
    # 1) 셀 텍스트
    try:
        wb = load_workbook(io.BytesIO(raw), data_only=True)
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                vals = [str(c) for c in row if c is not None]
                if vals:
                    text_parts.append("\t".join(vals))
    except Exception as e:
        print("xlsx text parse err:", e)
    # 2) 내장 이미지 (ZIP 스캔)
    try:
        with zipfile.ZipFile(io.BytesIO(raw)) as zf:
            for name in zf.namelist():
                low = name.lower()
                if low.startswith("xl/media/") and low.split(".")[-1] in (
                    "png", "jpg", "jpeg", "bmp", "gif", "tif", "tiff", "webp"
                ):
                    ocr_parts.append(ocr_image_bytes(zf.read(name)))
    except Exception as e:
        print("xlsx media zip scan err:", e)
    return "\n".join(text_parts).strip(), [t for t in ocr_parts if t]


def try_guess_image_from_bin(bin_bytes: bytes) -> bytes:
    """HWP OLE BinData 스트림에서 PNG/JPEG 시그니처를 찾아 잘라내기"""
    if not bin_bytes:
        return b""
    sigs = [b"\x89PNG\r\n\x1a\n", b"\xff\xd8\xff"]  # PNG/JPEG
    for sig in sigs:
        idx = bin_bytes.find(sig)
        if idx != -1:
            return bin_bytes[idx:]
    return b""

def extract_from_hwp_images_only(raw: bytes) -> Tuple[str, List[str]]:
    """
    HWP (OLE Compound)에서 BinData류 스트림만 뒤져서 이미지 OCR.
    본문 텍스트는 별도 파서(hwp5txt)가 없으면 생략.
    """
    ocr_parts = []
    with tempfile.NamedTemporaryFile(delete=False, suffix=".hwp") as tmp:
        tmp.write(raw)
        tmp.flush()
        path = tmp.name
    try:
        if not olefile.isOleFile(path):
            return "", []
        with olefile.OleFileIO(path) as ole:
            for entry in ole.listdir(streams=True, storages=True):
                # 예: ['BinData', 'Bin0001'] 등
                if any("Bin" in e or "BIN" in e for e in entry):
                    try:
                        stream = ole.openstream(entry).read()
                        img = try_guess_image_from_bin(stream)
                        if img:
                            ocr_parts.append(ocr_image_bytes(img))
                    except Exception as e:
                        print("hwp bin read err:", e)
    finally:
        try:
            os.remove(path)
        except:
            pass
    return "", [t for t in ocr_parts if t]

def extract_text_from_hwp_hwp5txt(raw: bytes) -> str:
    """
    (선택) hwp5txt CLI가 설치된 환경에서 HWP 본문 텍스트까지 추출.
    설치가 안돼 있으면 빈 문자열 반환.
    """
    with tempfile.NamedTemporaryFile(delete=False, suffix=".hwp") as tmp:
        tmp.write(raw)
        tmp.flush()
        path = tmp.name
    try:
        # hwp5txt 가 PATH 에 있어야 함 (리눅스 환경 권장)
        out = subprocess.check_output(
            ["hwp5txt", path], encoding="utf-8", errors="ignore")
        return out.strip()
    except Exception as e:
        print("hwp5txt not available or failed:", e)
        return ""
    finally:
        try:
            os.remove(path)
        except:
            pass

def extract_from_hwpx_zip(raw: bytes) -> Tuple[str, List[str]]:
    import zipfile
    import io
    text_parts, ocr_parts = [], []
    try:
        with zipfile.ZipFile(io.BytesIO(raw)) as zf:
            for name in zf.namelist():
                low = name.lower()
                # HWPX는 주로 Contents/Resources/… 경로에 이미지가 있음
                if ("/resources/" in low or low.startswith("contents/")) and low.split(".")[-1] in (
                    "png", "jpg", "jpeg", "bmp", "gif", "tif", "tiff", "webp"
                ):
                    ocr_parts.append(ocr_image_bytes(zf.read(name)))
            # (원하면 여기서 XML 텍스트도 추가 파싱)
    except Exception as e:
        print("hwpx zip scan err:", e)
    return "\n".join(text_parts).strip(), [t for t in ocr_parts if t]
# ---------------------------
# 메인 디스패처
# ---------------------------
def extract_all_text_and_images(binary: bytes, filename: str) -> str:
    """
    파일 확장자에 따라 본문 텍스트 + 내장 이미지 OCR을 합쳐 하나의 문자열로 반환.
    프론트 수정 없이 `text` 하나로 내려주기 위함.
    """
    ext = (filename or "").lower().split(".")[-1]
    body, ocr_list = "", []

    try:
        if ext in ["png", "jpg", "jpeg", "bmp", "gif", "tif", "tiff", "webp"]:
            body, ocr_list = extract_from_image(binary)

        elif ext == "pdf":
            body, ocr_list = extract_from_pdf(binary)

        elif ext == "docx":
            body, ocr_list = extract_from_docx(binary)

        elif ext in ["pptx", "ppt"]:
            body, ocr_list = extract_from_pptx(binary)

        elif ext in ["xlsx", "xls"]:
            body, ocr_list = extract_from_xlsx(binary)

        elif ext == "txt":
            body, _ = extract_from_txt(binary)

        elif ext == "hwpx":
            body, ocr_list = extract_from_hwpx_zip(binary)

        elif ext == "hwp":

            body = extract_text_from_hwp_hwp5txt(binary)

            _, ocr_list = extract_from_hwp_images_only(binary)

        else:

            body = detect_text_encoding_and_decode(binary)

    except Exception as e:

        print("parse error:", type(e).__name__, e)

    # 최종 합치기
    body = (body or "").strip()
    ocr_text = "\n\n".join([t for t in (ocr_list or []) if t]).strip()

    if body and ocr_text:
        return f"{body}\n\n{ocr_text}"
    elif ocr_text:
        return ocr_text
    else:
        return body

def detect_text_encoding_and_decode(raw: bytes) -> str:
    """
    bytes → str 디코딩. chardet로 인코딩 추정, 실패 시 안전한 폴백.
    """
    if not raw:
        return ""
    try:
        import chardet
        guess = chardet.detect(raw) or {}
        enc = (guess.get("encoding") or "utf-8").strip()

        return raw.decode(enc, errors="replace")
    except Exception:
        try:
            return raw.decode("utf-8", errors="replace")
        except Exception:
            return raw.decode("latin-1", errors="replace")


load_dotenv()

app = FastAPI()

ALLOW_ORIGINS = [
    "http://127.0.0.1:5500",
    "http://localhost:5500",
    "http://127.0.0.1:3000",
    "http://localhost:3000",
    "http://127.0.0.1:8000",
    "http://localhost:8000",
    "https://hj-sp.github.io",
    "https://crystal-0109.github.io",
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=ALLOW_ORIGINS,
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=["*"],
    max_age=600,
)

@app.get("/whoami")
def whoami():
    return {
        "ok": True,
        "has_json_vision": bool(os.environ.get("GOOGLE_APPLICATION_CREDENTIALS_JSON")),
        "has_json_translate": bool(os.environ.get("GOOGLE_CREDENTIALS_JSON")),
        "has_path_var": bool(os.environ.get("GOOGLE_APPLICATION_CREDENTIALS")),
        "routes": [r.path for r in app.routes],
    }

@app.get("/cors-test")
def cors_test():
    return {"ok": True}


# ---- Vision (GOOGLE_APPLICATION_CREDENTIALS_JSON) ----
try:
    _vision_info = json.loads(
        os.environ["GOOGLE_APPLICATION_CREDENTIALS_JSON"])
    _vision_creds = service_account.Credentials.from_service_account_info(
        _vision_info)
    vision_client = vision.ImageAnnotatorClient(credentials=_vision_creds)
except Exception as e:
    # 초기화 실패 시 디버그 로그
    print("[Vision Init Error]", e)
    vision_client = None

# ---- Translate (GOOGLE_CREDENTIALS_JSON) ----
try:
    _trans_info = json.loads(os.environ["GOOGLE_CREDENTIALS_JSON"])
    _trans_creds = service_account.Credentials.from_service_account_info(
        _trans_info)
    translate_client = google_translate.Client(credentials=_trans_creds)
except Exception as e:
    print("[Translate Init Error]", e)
    translate_client = None

class TextInput(BaseModel):
    content: str
    style: str = "default"
    offset: int = 0
    source: str = "rewrite"

class StyleChangeRequest(BaseModel):
    text: str
    style: str


AudioSegment.converter = imageio_ffmpeg.get_ffmpeg_exe()  # 패키지 내부 ffmpeg 사용

MISTRAL_API_KEY_H = os.getenv("MISTRAL_API_KEY_H")
MISTRAL_API_KEY_S = os.getenv("MISTRAL_API_KEY_S")
AGENT_ID_SUMMARY = os.getenv("MISTRAL_AGENT_ID_SUMMARY")
AGENT_ID_REWRITE = os.getenv("MISTRAL_AGENT_ID_REWRITE")
AGENT_ID_GRAMMAR = os.getenv("MISTRAL_AGENT_ID_GRAMMAR")
AGENT_ID_GRAMMAR2 = os.getenv("MISTRAL_AGENT_ID_GRAMMAR2")
AGENT_ID_GRAMMAR3 = os.getenv("MISTRAL_AGENT_ID_GRAMMAR3")
COHERE_API_KEY = os.getenv("COHERE_API_KEY")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
client = OpenAI(api_key=OPENAI_API_KEY)
PAPAGO_CLIENT_ID = os.getenv("PAPAGO_CLIENT_ID")
PAPAGO_CLIENT_SECRET = os.getenv("PAPAGO_CLIENT_SECRET")

@app.post("/searchExample")
async def search_example(data: TextInput):
    count = 5  # 고정
    prompt = f"""
    '{data.content}'이라는 표현을 다양한 문맥과 상황 속에서 자연스럽게 사용하는 예문을 {count}개 만들어줘.
    {'앞에서 제공된 예문과 겹치지 않게' if data.offset > 0 else ''}
    각 예문은 문맥이 풍부해야 하며, 문장 길이는 30자에서 100자 사이로 다양하게 작성해줘.
    모든 예문은 번호를 매기고, 번호별로 줄바꿈해서 구분해줘.
    """

    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "너는 국어 예문 생성 도우미야."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7
        )
        output = response.choices[0].message.content
        return {"result": output}
    except Exception as e:
        return {"error": f"GPT 호출 오류: {str(e)}"}


@app.post("/gptStyleChange")
async def gpt_style_change(request: Request):
    body = await request.json()

    text = body.get('text')
    style = body.get('style')

    if not text or not style:
        print("❗ text나 style이 비어있음")
        return {"error": "text 또는 style이 비어있습니다"}

    style = style.lower()

    style_instructions = {
        "formal": "문장을 격식 있고 정중하게 바꿔줘.",
        "casual": "문장을 친근하고 일상적인 말투로 바꿔줘.",
        "literary": "문장을 문학적으로 품격 있게 다시 써줘.",
        "academic": "문장을 전문적이고 학술적인 표현으로 바꿔줘."
    }

    instruction = style_instructions.get(style)

    if instruction is None:
        print(f"❗ 지원하지 않는 스타일: {style}")
        return {"error": f"지원하지 않는 스타일입니다. (서버가 받은 style: {style})"}

    full_prompt = f"아래 문장을 {instruction}\n\n'{text}'"

    try:
        completion = client.chat.completions.create(
            model="gpt-4",
            messages=[
                {"role": "system",
                    "content": "너는 문체를 변환하는 전문 AI 어시스턴트야. 문체 변경시 비속어는 사용하지 마."},
                {"role": "user", "content": full_prompt}
            ],
            temperature=0.7
        )
        output = completion.choices[0].message.content
        print(output)

        # 글 앞 뒤 따옴표 제거
        output = output[1:-1] if len(
            output) >= 2 and output[0] == output[-1] and output[0] in ('"', "'") else output
        print(output)

        return {"styled_text": output}
    except Exception as e:
        return {"error": f"GPT 호출 오류: {str(e)}"}

@app.post("/mistralRewrite")
async def mistral_rewrite(content: TextInput):
    headers = {
        "Authorization": f"Bearer {MISTRAL_API_KEY_H}",
        "Content-Type": "application/json"
    }
    count = 1

    prompt = f"""너는 창의적이고 유연한 글쓰기 첨삭 도우미야. 아래 글을 바탕으로 리라이팅한 **예시문을 반드시 하나만** 작성해 줘. 다음 지침을 반드시 따라야 해:

1. 문체와 의미는 유지하되,  
2. 문장 구조(예: 어순, 구문 유형, 능동/수동, 문장 길이)를 다양하게 바꾸고,  
3. 단어 선택과 어휘 스타일(예: 묘사 중심, 감정 강조, 간결체, 문어체 등)을 매번 다르게 써줘.

출력 형식은 다음과 같아:

(여기에 리라이팅된 문장)

아무 설명 없이 예시문 하나만 보여줘.
원문: {content.content}
"""

    payload = {
        "agent_id": AGENT_ID_REWRITE,
        "messages": [
            {"role": "user", "content": prompt}
        ]
    }

    response = requests.post(
        "https://api.mistral.ai/v1/agents/completions",
        headers=headers,
        json=payload
    )

    if response.status_code != 200:
        return {"error": f"HTTP 오류: {response.status_code}", "detail": response.text}

    result = response.json()

    try:
        message = result["choices"][0]["message"]["content"]
    except (KeyError, IndexError) as e:
        return {"error": "응답 파싱 오류", "detail": str(e), "raw_response": result}

    return {"result": message}

def strip_mode_banner(s: str) -> str:
    if not s:
        return s
    # 맨 앞줄에 오는 "• **요약 모드: basic**" 같은 변형들 일괄 제거
    s = re.sub(
        r'^\s*(?:[-*•]\s*)?(?:\*{1,3})?\s*요약\s*모드\s*:\s*[^\n]*\n+',
        '',
        s,
        flags=re.IGNORECASE
    )
    return s.lstrip()


@app.post("/summary")
async def summarize(content: TextInput):
    headers = {
        "Authorization": f"Bearer {MISTRAL_API_KEY_H}",
        "Content-Type": "application/json"
    }
    payload = {
        "agent_id": AGENT_ID_SUMMARY,
        "messages": [
            {
                "role": "user",
                "content": f"다음 글을 최대한 간결하게 요약해줘. 대신 핵심내용은 포함해줘:\n\n{content.content}"
            }
        ]
    }
    response = requests.post(
        "https://api.mistral.ai/v1/agents/completions",
        headers=headers, json=payload
    )
    result = response.json()
    try:
        message = result["choices"][0]["message"]["content"]
    except (KeyError, IndexError):
        return {"error": "Mistral 요약 오류"}
    # ✅ 라벨 제거
    return {"result": strip_mode_banner(message)}

@app.post("/expand")
async def expand(payload: ExpandInput):
    try:
        text = (payload.content or "").strip()
        if not text:
            raise HTTPException(status_code=400, detail="content is empty")

        # 1) 길이 증가 퍼센트 확정
        level_map = {'low': 20, 'medium': 50, 'high': 80, 'xhigh': 100}
        if getattr(payload, "length_level", None):
            boost_in = level_map.get(payload.length_level.lower(), 50)
        else:
            boost_in = int(getattr(payload, "length_boost", 20) or 20)

        allowed = [20, 50, 70, 80, 100]
        # 허용값에 스냅 → 미세값도 근접 허용으로 매핑
        boost = boost_in if boost_in in allowed else min(
            allowed, key=lambda x: abs(x - boost_in)
        )

        # 2) 입력 길이/문장수 측정 (kss)
        orig_chars = len(text)
        orig_sents = count_ko_sentences(text)
        short_input = (orig_chars < 120) or (orig_sents <= 2)

        # 3) 목표/최소/최대 글자수 설정
        target_chars = int(orig_chars * (1 + boost / 100.0))
        min_chars = int(orig_chars * (1 + (boost / 100.0) * 0.95))
        max_chars = int(orig_chars * (1 + boost / 100.0 + 0.08))

        if short_input:
            # 짧은 입력은 바닥/상한을 올려 자연스러운 분량 확보
            floor_min = orig_chars + (120 if boost >= 70 else 80)
            floor_target = floor_min + 60
            floor_max = floor_target + 160
            min_chars = max(min_chars, floor_min)
            target_chars = max(target_chars, floor_target)
            max_chars = max(max_chars, floor_max)

        # 4) 문장 추가 옵션
        add_n = max(0, min(int(getattr(payload, "add_sentences", 0) or 0), 50))

        # 짧은 입력 + add_n 미지정 시 기본 1~2문장 보강
        if short_input and add_n == 0:
            add_n = 2 if boost >= 70 else 1

        if add_n > 0:
            # 문장 추가 시 예상 글자수(문장당 70자 기준)로 예산 상향
            extra_chars = add_n * 70
            min_chars = max(min_chars, int(orig_chars + extra_chars * 0.8))
            target_chars = max(target_chars, min_chars + 60)
            max_chars = max(max_chars, int(orig_chars + extra_chars * 1.4))

        # 불변성 보장
        if min_chars > max_chars:
            max_chars = min_chars + 120
        if target_chars > max_chars:
            target_chars = max(min_chars, max_chars - 60)

        # 5) 1차 생성
        system_msg = "당신은 한국어 글을 자연스럽고 명확하게 확장하는 전문가입니다."
        guide = [
            "원문 의미를 유지하고 전개를 매끄럽게 확장하라.",
            "단순한 바꿔쓰기(패러프레이즈)만 하지 말고, 새로운 설명/이유/예시/전환을 덧붙여라.",
            f"최소 {min_chars}자 이상, 최대 {max_chars}자 이하를 엄수하라.",
            f"목표 길이: 약 {target_chars}자.",
            "불필요한 도입/결론/중복 요약을 피하라.",
            "모든 문장은 완결형(…다./…요./…입니다.)으로 끝내라.",
        ]
        if add_n > 0:
            per_sent_buf = 130  # 문장 하나당 안전 버퍼
            max_chars += add_n * per_sent_buf
            guide.append(f"원문 내용은 유지하되, **최소 {add_n}문장**을 새로 추가하라.")

        user_prompt = f"[원문]\n{text}\n\n[지시]\n- " + "\n- ".join(guide)

        out, finish = _expand_llm(
            client,
            [{"role": "system", "content": system_msg},
             {"role": "user", "content": user_prompt}],
            max_chars=max_chars,
            temperature=0.35
        )

        if len(out) > max_chars:
            out = _expand_smart_trim_to_chars(out, max_chars)

        # 6) 끝맺음 보정
        if not _expand_ends_with_terminal(out) or finish == "length":
            remain = max(0, max_chars - len(out))
            if remain > 24:
                tail_ctx = out[-500:] if len(out) > 500 else out
                tail_prompt = (
                    f"{tail_ctx}\n\n"
                    "위 글의 마지막을 자연스럽게 끝맺으세요.\n"
                    "- 한 문장만 출력\n"
                    "- 새로운 주장/예시 추가 금지, 앞 맥락만 완결\n"
                    f"- 최대 길이 {min(140, remain)}자\n"
                    "- 완결형 종결어미로 끝내세요(…다./…요./…입니다.)."
                )
                tail, _ = _expand_llm(
                    client,
                    [{"role": "system", "content": "간결한 문장 마무리 도우미"},
                     {"role": "user", "content": tail_prompt}],
                    max_chars=min(140, remain),
                    temperature=0.2
                )
                if tail:
                    sep = "" if out.endswith((" ", "\n")) else " "
                    out = (out + sep + tail).strip()
            if len(out) > max_chars:
                out = _expand_smart_trim_to_chars(out, max_chars)

        # 7) 최소 글자수 미달 시 보강
        tries = 0
        while len(out) < min_chars and tries < 2:
            need_chars = max(120, min(min_chars - len(out),
                             int((min_chars - len(out)) * 1.2)))
            add_guide = [
                "기존 문장은 바꾸지 말고, '추가 문장/문단'만 이어서 작성하라.",
                f"현재 길이: {len(out)}자, 목표 최소: {min_chars}자 → 대략 {need_chars}자 보강.",
                "중복 요약·도입/결론 재작성 금지.",
                "완결형 종결어미로 끝내라."
            ]
            cont_prompt = (
                f"[지금까지의 초안]\n{out}\n\n"
                "[지시]\n- " + "\n- ".join(add_guide) + "\n"
                "※ 출력은 '추가되는 내용만' 반환하세요."
            )
            remain_budget = max_chars - len(out)
            if remain_budget <= 60:
                break
            add_text, _ = _expand_llm(
                client,
                [{"role": "system", "content": "확장 보강 도우미"},
                 {"role": "user", "content": cont_prompt}],
                max_chars=min(remain_budget, max(220, need_chars + 140)),
                temperature=0.35
            )
            add_text = (add_text or "").strip()
            if len(add_text) > remain_budget:
                add_text = _expand_smart_trim_to_chars(add_text, remain_budget)
            if add_text:
                sep = "\n\n" if not out.endswith("\n") else "\n"
                out = (out + sep + add_text).strip()
            if len(out) > max_chars:
                out = _expand_smart_trim_to_chars(out, max_chars)
            if not _expand_ends_with_terminal(out):
                out = _expand_crop_to_last_boundary(out)
            tries += 1

        # 8) 문장 추가 정확도 루프 (kss 기반)
        if add_n > 0:
            per_sent_buf = 130  # 문장 하나당 안전 버퍼
            max_chars += add_n * per_sent_buf
            needed = (orig_sents + add_n) - count_ko_sentences(out)
            s_tries = 0
            while needed > 0 and s_tries < 4:
                remain_budget = max_chars - len(out)
                if remain_budget <= 80:
                    break

                per_sent_budget = 120
                ask_budget = min(remain_budget, needed * per_sent_budget + 40)

                cont_prompt = f"""
[지금까지의 초안]
{out}

[지시]
- 기존 문장은 수정하지 말고 '이어 쓰기'만 하세요.
- **정확히 {needed}개 문장**을 추가하세요.
- 번호/불릿/따옴표/머리말 없이, 각 문장은 **한 줄에 한 문장**만 출력하려 노력하세요.
- 모든 문장은 완결형(…다./…요./…입니다.)으로 끝내세요.
- 도입/결론/요약/중복 금지, 앞 맥락만 자연스럽게 보강하세요.
- 출력은 **추가되는 문장만** 반환하세요.
""".strip()

                add_text, _ = _expand_llm(
                    client,
                    [
                        {"role": "system", "content": "문장 추가 전용 보조"},
                        {"role": "user", "content": cont_prompt}
                    ],
                    max_chars=ask_budget,
                    temperature=0.3
                )
                add_text = (add_text or "").strip()
                if not add_text:
                    s_tries += 1
                    continue

                # (1) kss로 문장 단위 분리
                cand_sents = split_ko_sentences(add_text)

                # (2) 불릿/번호 제거 + 종결 보정
                cleaned = []
                for s in cand_sents:
                    s = re.sub(r'^[\s\-\*\•\·]+', '', s)      # 불릿류 제거
                    s = re.sub(r'^\s*\d+[\.\)]\s*', '', s)     # 번호 제거
                    s = s.strip()
                    if not s:
                        continue
                    if not s.endswith(("다.", "요.", "입니다.", "?", "!", "…", ".")):
                        s = s.rstrip(" ,;:") + "다."
                    cleaned.append(s)

                if not cleaned:
                    s_tries += 1
                    continue

                # (3) 필요한 개수만 선별
                to_take = cleaned[:needed]

                # (4) 예산 초과 시 뒤에서 문장 단위로 줄이기
                room = remain_budget - (0 if out.endswith("\n") else 1)
                while to_take and len("\n".join(to_take)) > room:
                    to_take.pop()

                if not to_take:
                    s_tries += 1
                    continue

                # (5) 본문에 붙이기
                sep = "\n" if out.endswith("\n") else "\n\n"
                out = (out + sep + "\n".join(to_take)).strip()

                if len(out) > max_chars:
                    out = _expand_smart_trim_to_chars(out, max_chars)

                # (6) 남은 필요 문장 재계산
                needed = (orig_sents + add_n) - count_ko_sentences(out)
                s_tries += 1

            # (옵션) 마지막 미달 문장 보정 1회
            final_needed = (orig_sents + add_n) - count_ko_sentences(out)
            remain_budget = max_chars - len(out)
            if final_needed > 0 and remain_budget > 60:
                per_sent_budget = 90
                ask_budget = min(
                    remain_budget, final_needed * per_sent_budget + 20)
                cont_prompt = f"""
[지금까지의 초안]
{out}

[지시]
- 기존 문장 수정 금지, 이어쓰기만.
- **정확히 {final_needed}개 문장** 추가.
- 번호/불릿 없이 자연스러운 완결형 문장.
- 출력은 추가 문장만.
""".strip()
                add_text, _ = _expand_llm(
                    client,
                    [
                        {"role": "system", "content": "문장 추가(마감) 보조"},
                        {"role": "user", "content": cont_prompt}
                    ],
                    max_chars=ask_budget,
                    temperature=0.25
                )
                lines = split_ko_sentences(add_text or "")
                take = lines[:final_needed]

                room = remain_budget - (0 if out.endswith("\n") else 1)
                while take and (len("\n".join(take)) > room):
                    take.pop()

                if take:
                    sep = "\n" if out.endswith("\n") else "\n\n"
                    out = (out + sep + "\n".join(take)).strip()
                    if len(out) > max_chars:
                        out = _expand_smart_trim_to_chars(out, max_chars)

        # 9) 최종 최댓값 보정
        if len(out) > max_chars:
            out = _expand_smart_trim_to_chars(out, max_chars)

        return {
            "result": out,
            "meta": {
                "orig_chars": orig_chars,
                "min_chars": min_chars,
                "max_chars": max_chars,
                "target_chars": target_chars,
                "requested_add_sentences": add_n,
                "final_chars": len(out),
                "final_sents": count_ko_sentences(out),
            }
        }

    except HTTPException:
        raise
    except Exception as e:
        return JSONResponse(
            status_code=500,
            content={"error": f"/expand failed: {type(e).__name__}: {str(e)}"}
        )

@app.post("/mistralGrammar")
async def mistral_grammar(content: TextInput):
    start_time = time.perf_counter()  # 처리 시작 시점
    headers = {
        "Authorization": f"Bearer {MISTRAL_API_KEY_S}",
        "Content-Type": "application/json"
    }
    payload = {
        "agent_id": AGENT_ID_GRAMMAR3,
        "messages": [
            {"role": "user", "content": content.content}
        ]
    }

    response = requests.post(
        "https://api.mistral.ai/v1/agents/completions",
        headers=headers,
        json=payload
    )

    if response.status_code != 200:
        return {"error": f"HTTP 오류: {response.status_code}", "detail": response.text}

    result = response.json()
    print("Mistral 응답:", result)

    try:
        message = result["choices"][0]["message"]["content"]
    except (KeyError, IndexError) as e:
        return {"error": "응답 파싱 오류", "detail": str(e), "raw_response": result}

    end_time = time.perf_counter()
    elapsed_time = end_time - start_time
    print(f"[총 처리 시간] {elapsed_time:.2f}초")
    print("현재 시간 : ", datetime.now().time())

    return {"result": message}

def split_sentences(text):
    text = text.strip()
    if not text:
        return []

    # 문장 끝 구두점이 하나라도 포함돼 있다면 그 기준으로 분리
    if re.search(r'[.!?]', text):
        # 마침표/느낌표/물음표 기준으로 문장 분리 (캡처 그룹으로 포함시켜서 문장 구두점도 살림)
        parts = re.split(r'([.!?])', text)
        sentences = []
        for i in range(0, len(parts) - 1, 2):
            sentence = parts[i].strip() + parts[i + 1]  # 문장 + 구두점
            if sentence.strip():
                sentences.append(sentence.strip())
        return sentences
    else:
        # 마침표/느낌표/물음표 없으면 한 문장으로 간주
        return [text]

# 기존
# @app.post("/cohereHonorific")
# async def cohere_honorific(content: TextInput):
#     co = cohere.ClientV2(COHERE_API_KEY)
#     user_prompt = f"{content.content}\n\n이 글에서 ~습니다. 처럼 높임말로 바꿔줘. 그리고 결과는 딱 글만 나오게 해줘."
#     response = co.chat(
#         model="command-a-03-2025",
#         messages=[{"role": "user", "content": user_prompt}]
#     )
#     full_text = response.message.content[0].text
#     return {"result": full_text}


@app.post("/cohereHonorific")
async def cohere_honorific(request: Request):
    body = await request.json()
    text = (body.get("content") or body.get("text") or "").strip()
    level = (body.get("level") or "hamnida").lower()
    if level not in ("haeyo", "hamnida"):
        level = "hamnida"

    if not text:
        return {"error": "content(text)가 비어있습니다."}

    style_map = {
        "haeyo": "문장을 정중하지만 덜 격식 있는 ‘해요체(~요)’로 자연스럽게 바꿔줘.",
        "hamnida": "문장을 공식/문어체의 ‘합니다체(~습니다/~입니다)’로 자연스럽게 바꿔줘.",
    }
    instruction = style_map.get(level, style_map["hamnida"])

    user_prompt = (
        f"{text}\n\n"
        f"{instruction}\n"
        f"출력은 변환된 본문만, 설명/따옴표/머리말 없이 그대로 내보내."
    )

    co = cohere.ClientV2(COHERE_API_KEY)
    response = co.chat(
        model="command-a-03-2025",
        messages=[{"role": "user", "content": user_prompt}]
    )
    out = response.message.content[0].text
    return {"result": out, "level": level}

@app.post("/cohereInformal")
async def cohere_informal(payload: InformalInput):
    import re
    co = cohere.ClientV2(COHERE_API_KEY)

    # 0) 입력
    src = (payload.content or "").strip()
    if not src:
        return {"error": "content가 비었습니다."}

    # 1) 엔딩 옵션
    ending = (payload.ending or "").strip().lower()
    ending_rules = {
        "hada": "모든 문장을 '~다'로 끝나는 평서형(…한다/…이다)으로 바꿔줘. 구어체/해요체/존칭 표현은 쓰지 마.",
        "haetda": "모든 문장을 과거 평서형 '~했다/~였다'로 바꿔줘. 구어체/해요체/존칭 표현은 쓰지 마.",
        "hae": "모든 문장을 '~해/~했어' 반말체로 바꿔줘. 머리말이나 설명을 붙이지 마."
    }
    instruction = ending_rules.get(ending, ending_rules["hae"])

    # 2) 태그+few-shot 예시: 초단문/단문에서도 본문만 나오게
    user_prompt = f"""
다음 규칙을 지켜서 문장을 변환해.
- {instruction}
- 결과는 반드시 <out>와 </out> 사이에만 작성해.
- 설명/머리말/따옴표/불릿/코드블록/제목 금지. 오직 변환 결과만.
- </out> 뒤에는 아무 것도 쓰지 마.

[예시]
<in>안녕하세요</in>
<out>안녕</out>

<in>감사합니다</in>
<out>고마워</out>

<in>죄송합니다</in>
<out>미안해</out>

<in>네</in>
<out>응</out>

[변환할 문장]
<in>{src}</in>

[출력]
<out>"""

    resp = co.chat(
        model="command-a-03-2025",
        messages=[
            {"role": "system", "content": "너는 한국어 문체 변환기야. 오직 결과만 출력해."},
            {"role": "user", "content": user_prompt}
        ],
        stop_sequences=["</out>"]  # 태그 닫히면 즉시 중단
    )

    raw = (resp.message.content[0].text or "").strip()

    # 3) 파싱: <out> ... </out> (</out>은 stop으로 잘렸을 수 있음)
    m = re.search(r"<out>([\s\S]*?)$", raw)
    out = (m.group(1) if m else raw).strip()

    # 4) 서버측 살균기: 머리말/메타 라인/코드펜스/따옴표 제거
    def sanitize(s: str) -> str:
        # 코드펜스
        s = re.sub(r"^```[\s\S]*?```$", "", s, flags=re.MULTILINE).strip()
        # 양끝 따옴표
        if len(s) >= 2 and s[0] in "'\"" and s[-1] == s[0]:
            s = s[1:-1].strip()
        # 흔한 메타 헤더 라인 제거
        ban_prefix = r"(답변이|반말|높임말|알겠어|설명|지시|출력|예시|머리말)"
        lines = [ln for ln in s.splitlines()
                 if not re.match(rf"^\s*{ban_prefix}\b", ln)]
        s = "\n".join(lines).strip()
        # 남은 줄에서 헤더 비슷한 것(짧은 한 단어) 제거
        lines = [ln for ln in s.splitlines() if not re.fullmatch(
            r"\s*[가-힣A-Za-z]{1,6}\s*", ln)]
        s = "\n".join(lines).strip()
        return s

    out = sanitize(out)
    return {"result": out, "ending": ending or "hae"}


@app.post("/translate")
async def translate_text(request: Request):
    body = await request.json()
    text = body.get("text")
    source = body.get("source", "None")
    target = body.get("target", "en")

    if not text:
        return {"error": "No text provided"}

    if translate_client is None:
        return {"error": "Translate client not initialized. Check GOOGLE_CREDENTIALS_JSON."}

    try:
        if source and source != "auto":
            result = translate_client.translate(
                text, source_language=source, target_language=target
            )
        else:
            result = translate_client.translate(text, target_language=target)

        translated_clean = html.unescape(result["translatedText"])
        return {"result": translated_clean}
    except Exception as e:
        return {"error": f"Google 번역 API 호출 오류: {str(e)}"}

def extract_pdf_text(pdf_path: str) -> str:
    """
    PyPDF2로 PDF의 전체 텍스트를 추출하고, 온점 앞 공백 제거 등 간단 후처리.
    """
    try:
        reader = PdfReader(pdf_path)
        extracted_text = ""
        for page in reader.pages:
            extracted_text += page.extract_text() or ""
        # 온점 앞의 띄어쓰기 제거
        cleaned_text = re.sub(r" (?=\.)", "", extracted_text)
        return cleaned_text.strip()
    except Exception as e:
        raise RuntimeError(f"PDF 처리 중 오류: {e}")


def libreoffice_to_pdf(input_path: str) -> str:
    """
    LibreOffice(headless)로 거의 모든 오피스/HWP 포맷을 PDF로 변환.
    예: soffice --headless --convert-to pdf --outdir /tmp input.hwp
    """
    outdir = tempfile.gettempdir()

    cmd = ["soffice", "--headless", "--convert-to",
           "pdf", "--outdir", outdir, input_path]
    try:
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE,
                       stderr=subprocess.PIPE)
        stem = pathlib.Path(input_path).stem
        out_pdf = os.path.join(outdir, f"{stem}.pdf")
        if not os.path.exists(out_pdf):
            raise RuntimeError("LibreOffice 변환 결과 PDF를 찾을 수 없습니다.")
        return out_pdf
    except subprocess.CalledProcessError as e:
        raise RuntimeError(
            f"LibreOffice 변환 실패: {e.stderr.decode('utf-8', errors='ignore')}")

def extract_hwpx_text(local_path: str) -> str:
    """
    HWPX는 OOXML 유사 구조의 '압축(zip)+XML' 포맷.
    Contents/*.xml 에서 모든 텍스트 노드({*}t)를 긁어와 문단 단위로 합칩니다.
    """
    texts = []
    with zipfile.ZipFile(local_path, "r") as zf:

        xml_names = [n for n in zf.namelist()
                     if n.lower().startswith("contents/") and n.lower().endswith(".xml")]
        if not xml_names:

            xml_names = [n for n in zf.namelist(
            ) if n.lower().endswith(".xml")]

        for name in sorted(xml_names):
            with zf.open(name) as fp:
                data = fp.read()
            try:
                root = ET.fromstring(data)
            except ET.ParseError:
                continue

            for node in root.findall(".//{*}t"):
                if node.text and node.text.strip():
                    texts.append(node.text)

    text = "\n".join(texts)
    text = re.sub(r"[ \t]+\n", "\n", text)         #
    text = re.sub(r"\s{2,}", " ", text)            #
    text = re.sub(r"\s+(?=[\.\,\!\?\:\;])", "", text)
    return text.strip()

def extract_text_by_ext(local_path: str, filename: str) -> str:
    """
    파일 확장자/포맷에 따라 텍스트를 추출.
    - 직접 추출 가능한 포맷: pdf, docx, pptx, xlsx, txt
    - 그 외(ppt, xls, hwp, hwpx, doc, rtf 등)는 LibreOffice로 pdf 변환 후 PDF 추출 재사용
    """
    ext = pathlib.Path(filename).suffix.lower().lstrip(".")
    text = ""

    if ext == "pdf":
        text = extract_pdf_text(local_path)

    elif ext == "docx":
        # mammoth: docx -> plain text
        try:
            import mammoth
            with open(local_path, "rb") as f:
                result = mammoth.extract_raw_text(f)
            text = (result.value or "").strip()
        except Exception as e:
            raise RuntimeError(f"DOCX 처리 오류: {e}")

    elif ext == "pptx":
        # python-pptx: 모든 슬라이드에서 shape.text 모으기
        try:
            from pptx import Presentation
            prs = Presentation(local_path)
            chunks = []
            for slide in prs.slides:
                for shp in slide.shapes:
                    if hasattr(shp, "text") and shp.text:
                        chunks.append(shp.text)
                # 노트 영역까지 필요하면 아래 활성화
                if getattr(slide, "notes_slide", None) and slide.notes_slide.notes_text_frame:
                    chunks.append(slide.notes_slide.notes_text_frame.text)
            text = "\n".join(chunks).strip()
        except Exception as e:
            raise RuntimeError(f"PPTX 처리 오류: {e}")

    elif ext == "xlsx":
        # openpyxl: 셀 텍스트를 시트별로 모으기
        try:
            import openpyxl
            wb = openpyxl.load_workbook(
                local_path, data_only=True, read_only=True)
            rows = []
            for ws in wb.worksheets:
                rows.append(f"### 시트: {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    # 너무 긴 엑셀이면 여기서 줄수/열수 제한 가능
                    rows.append("\t".join(cells))
            text = "\n".join(rows).strip()
        except Exception as e:
            raise RuntimeError(f"XLSX 처리 오류: {e}")

    elif ext == "txt":
        # 단순 텍스트
        try:
            with open(local_path, "rb") as f:
                raw = f.read()
            text = raw.decode("utf-8", errors="ignore").strip()
        except Exception as e:
            raise RuntimeError(f"TXT 처리 오류: {e}")

    elif ext == "hwpx":
        # ✅ LibreOffice 없이 바로 파싱
        text = extract_hwpx_text(local_path)

    elif ext in ("ppt", "xls", "hwp", "doc", "rtf"):
        # 이들은 계속 LibreOffice 변환을 사용
        pdf_path = libreoffice_to_pdf(local_path)
        text = extract_pdf_text(pdf_path)

    else:

        try:
            pdf_path = libreoffice_to_pdf(local_path)
            text = extract_pdf_text(pdf_path)
        except Exception as e:
            raise RuntimeError(f"LibreOffice 변환/추출 실패: {e}")

    return text

@app.post("/fileScan")
async def file_scan(file: UploadFile = File(...)):
    """
    하나의 업로드 엔드포인트로:
    - 모든 문서 포맷의 본문 텍스트 추출
    - 문서 내 '삽입 이미지'들 OCR 결과까지 병합
    반환: {"filename", "text"}  (프론트 수정 불필요)
    """
    suffix = pathlib.Path(file.filename).suffix
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        binary = await file.read()
        tmp.write(binary)
        tmp_path = tmp.name

    try:
        merged = extract_all_text_and_images(binary, file.filename)

        MAX_CHARS = 100_000
        if len(merged) > MAX_CHARS:
            merged = merged[:MAX_CHARS]

        return {"filename": file.filename, "text": merged}
    except Exception as e:
        return JSONResponse(status_code=400, content={"error": str(e)})
    finally:
        try:
            os.remove(tmp_path)
        except:
            pass


@app.post("/pdfScan")
async def upload_pdf(pdf: UploadFile = File(...)):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
        contents = await pdf.read()
        temp_pdf.write(contents)
        temp_pdf_path = temp_pdf.name

    try:
        cleaned_text = extract_pdf_text(temp_pdf_path)
        return {"filename": pdf.filename, "text": cleaned_text}
    except Exception as e:
        return {"error": f"PDF 처리 중 오류 발생: {str(e)}"}
    finally:
        try:
            os.remove(temp_pdf_path)
        except:
            pass

# 단독 자모 제거(후처리)
def clean_ocr_text(text: str) -> str:
    """
    단독 자모(ㅏ, ㄱ 등)나 잡음을 제거함.
    """
    cleaned_lines = []
    for line in text.split('\n'):
        stripped = line.strip()
        # 단독 자모나 1~2글자 잡음 제거
        if re.fullmatch(r'[ㄱ-ㅎㅏ-ㅣ\s]*', stripped):
            continue
        if len(stripped) <= 2 and re.search(r'[ㄱ-ㅎㅏ-ㅣ]', stripped):
            continue
        cleaned_lines.append(line)
    return '\n'.join(cleaned_lines)

@app.post("/visionOCR")
async def vision_ocr(image: UploadFile = File(...)):
    try:
        if vision_client is None:
            return {"result": "Vision client not initialized. Check GOOGLE_APPLICATION_CREDENTIALS_JSON."}

        contents = await image.read()
        image_content = vision.Image(content=contents)
        response = vision_client.text_detection(image=image_content)

        texts = response.text_annotations
        if not texts:
            return {"result": "텍스트를 찾을 수 없습니다."}

        raw_text = texts[0].description
        cleaned_text = clean_ocr_text(raw_text)
        return {"result": cleaned_text}
    except Exception as e:
        print("❌ 서버 에러:", e)
        return {"result": f"서버 에러: {e}"}

@app.post("/speech")
async def upload_audio(audio: UploadFile = File(...)):
    start_time = time.time()

    # 업로드된 오디오 파일 읽기
    audio_bytes = await audio.read()
    audio_segment = AudioSegment.from_file(BytesIO(audio_bytes), format="wav")

    # (선택) 오디오 전처리
    silence = AudioSegment.silent(duration=1000)
    audio_segment = silence + audio_segment
    audio_segment = audio_segment.set_frame_rate(
        16000).set_channels(1).set_sample_width(2)
    audio_segment = audio_segment.normalize()
    audio_segment += 6

    print("현재 평균 데시벨:", audio_segment.dBFS)

    # 전체 오디오 그대로 사용
    recognizer = sr.Recognizer()
    with BytesIO() as wav_buffer:
        audio_segment.export(wav_buffer, format="wav")
        wav_buffer.seek(0)
        with sr.AudioFile(wav_buffer) as source:
            # (선택) 잡음 보정
            recognizer.adjust_for_ambient_noise(source, duration=1)
            audio_data = recognizer.record(source)  # 전체 파일 한 번에 읽음

            try:
                result = recognizer.recognize_google(
                    audio_data, language="ko-KR", show_all=True)
                if "alternative" in result:
                    text = result["alternative"][0]["transcript"]
                else:
                    text = ""
            except sr.UnknownValueError:
                text = ""
                print("오디오 인식 실패")
            except sr.RequestError as e:
                text = ""
                print("Google Web Speech API 요청 오류:", e)

    elapsed_time = time.time() - start_time

    if (text):
        print("\n" + text + "\n")
    print(f"걸린 시간: {elapsed_time:.3f}초")
    print()

    return {"text": text, "time": round(elapsed_time, 3)}

@app.post("/editorGrammar")
async def editorGrammar(content: TextInput):
    print(content.content)
    print()

    result = spell_checker.check(content.content)

    print(f"원문        : {result.original}")
    print(f"수정된 문장 : {result.checked}")
    print(f"오류 수     : {result.errors}")
    print(f"단어별 상태 : {list(result.words.keys())}")
    print(f"검사 시간   : {result.time:.4f}초")
    print("-" * 40)

    headers = {
        "Authorization": f"Bearer {MISTRAL_API_KEY_S}",
        "Content-Type": "application/json"
    }
    payload = {
        "agent_id": AGENT_ID_GRAMMAR2,
        "messages": [
            {"role": "user", "content": result.checked}
        ]
    }
    response = requests.post(
        "https://api.mistral.ai/v1/agents/completions",
        headers=headers,
        json=payload
    )
    if response.status_code != 200:
        return {"error": f"HTTP 오류: {response.status_code}", "detail": response.text}

    text = response.json()

    try:
        message = text["choices"][0]["message"]["content"]
    except (KeyError, IndexError) as e:
        return {"error": "응답 파싱 오류", "detail": str(e), "raw_response": text}

    return {
        "original": result.original,
        # "checked": result.checked,
        "errors": result.errors,
        # "words": list(result.words.keys()),
        "time": result.time,
        "checked": message,
    }

@app.post("/promptChange")
async def prompt_change(request: Request):
    body = await request.json()
    content = (body.get('content') or "").strip()
    prompt  = (body.get('prompt')  or "").strip()
    norm    = re.sub(r"\s+", " ", prompt).lower()  # 공백/대소문자 정규화

    # 1) 예시 사전 (우선순위 순서대로 나열)
    EXAMPLES = [
        # 교수님 메일
        (r"(교수님|교수|prof|이메일|메일)", """제목: 상명대학교 교육혁신추진팀에서 진행하는 핵심역량 진단 시행 안내

안녕하십니까 교수님,

저는 상명대학교 교육혁신추진팀(SM-IN)에서 근무하고 있는 교육혁신관리자입니다. 이번에 전달드리는 내용은 우리 대학의 교육 혁신을 위해 SM-IN 팀이 계획한 핵심역량 진단에 대한 정보입니다.

진단은 상명대학교 전체 재학생을 대상으로 2025년 9월 4일부터 10월 10일까지 진행될 예정이며, 모든 학생들의 참여가 필요합니다. 진단에 참여하려면 샘물통합정보시스템에 접속 후 '학생기본'>'핵심역량진단'>'역량진단평가' 메뉴를 통해 접근할 수 있습니다. 진단 결과는 익일 확인 가능하며, 결과 확인 경로 또한 동일한 경로에서 진행됩니다. 

본 진단의 목적은 학생들이 핵심역량(전문지식탐구, 창의적문제해결, 융복합, 다양성존중, 윤리실천)을 어느 정도 함양하고 있는지 진단하고, 향후 교육 방향을 결정하는 데 도움을 주는 것입니다. 재학 기간 동안 SM-IN 핵심역량을 우수하게 함양한 학생은 최우수인증자로서 졸업 시 총장 명의의 상장(인중서)을 받을 수 있습니다.

그러니 매년 1회씩 진단에 참여하여 자신의 역량점수 변화를 확인하시기 바랍니다. 혹시 관련하여 문의사항이 있다면 (서울) 02-2287-6456, (천안) 041-550-5508로 연락주시길 바랍니다.

학생들의 학업진행에 이 글이 많은 도움이 될 수 있기를 바랍니다.

감사합니다. 

상명대학교 교육혁신추진팀
"""),
        # 가정통신문
        (r"(가정통신문|가정 통신문|가정-통신문)", """가 정 통 신 문

제목: 추운 겨울, 건강한 생활을 위한 운동의 중요성

학부모님 안녕하십니까?
어느덧 찬바람이 불고 기온이 내려가는 겨울이 찾아왔습니다.
추운 날씨로 인해 야외활동이 줄어드는 시기이지만, 규칙적인 운동 습관은 여전히 우리 아이들의 건강을 지키는 데 매우 중요합니다.

운동은 신체를 튼튼하게 할 뿐 아니라 정신적인 안정과 긍정적인 에너지를 유지하는 데에도 큰 도움이 됩니다. 특히 가벼운 실내 유산소 운동이나 스트레칭은 추운 날씨 속에서도 쉽게 실천할 수 있으며, 스트레스 완화, 면역력 강화, 수면의 질 향상 등에 효과적입니다.

하루 30분 정도의 가벼운 운동만으로도 충분히 건강을 지킬 수 있습니다.
무리하지 않고 자신의 체력에 맞는 수준으로 꾸준히 실천하는 것이 가장 중요합니다. 가족이 함께 실내에서 몸을 풀거나, 주말에 가벼운 산책을 하는 등 운동을 일상의 즐거운 습관으로 만들어 주시길 바랍니다.

추운 계절일수록 몸과 마음이 움츠러들기 쉽습니다.
가정에서도 자녀가 활기차고 건강한 겨울을 보낼 수 있도록 따뜻한 관심과 격려 부탁드립니다.

감사합니다.
○○초등학교장 (인)
"""),
    ]

    # 2) 예시 매칭 (가장 먼저 일치하는 것을 반환)
    for pattern, example_text in EXAMPLES:
        if re.search(pattern, norm):
            print(f"[promptChange] DEMO HIT → pattern={pattern}, prompt='{prompt[:80]}...'")
            return {"result": example_text}

    # 3) 예시가 아니면 GPT 호출 (기존 동작)
    print(f"[promptChange] GPT CALL → prompt='{prompt[:80]}...'")
    response = client.chat.completions.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "당신은 글을 수정해주는 작문 보조 도구입니다."},
            {"role": "user", "content": f"{content}위의 글을 아래 프롬프트에 맞게 수정해줘.\n\n{prompt}"}
        ]
    )
    message = response.choices[0].message.content
    return {"result": message}


@app.post("/promptAdd")
async def expand(request: Request):
    body = await request.json()

    before = body.get('before')
    after = body.get('after')
    prompt = body.get('prompt')

    response = client.chat.completions.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "당신은 글을 수정해주는 작문 보조 도구입니다."},
            {"role": "user", "content": f"현재 커서 전의 글: {before}\n\n현재 커서 이후의 글: {after}\n\n현재 커서를 기준으로 분리한 글들을 아래 프롬프트에 맞게 수정해줘. 결과는 추가할 글만 알려줘.\n\n{prompt}"}
        ]
    )
    message = response.choices[0].message.content

    # 글 앞 뒤 따옴표 제거
    message = message[1:-1] if len(
        message) >= 2 and message[0] == message[-1] and message[0] in ('"', "'") else message
    print(message)
    return {"result": message}
